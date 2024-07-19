import os
from pptx import Presentation
from PIL import Image
import io

class PPTReader:
    def __init__(self, file_path):
        self.prs = Presentation(file_path)

    def extract_content(self):
        slides = []
        for slide in self.prs.slides:
            slide_content = {'text': '', 'images': []}
            
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_content['text'] += shape.text + '\n'
                
                if shape.shape_type == 13:  # 13 is the enum for picture
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image = Image.open(io.BytesIO(image_bytes))
                        
                        # Save image to a temporary file
                        temp_image_path = "temp_image_{0}.png".format(len(slide_content['images']))
                        image.save(temp_image_path)
                        slide_content['images'].append(temp_image_path)
                    except AttributeError:
                        print("Warning: Failed to process an image in the slide. Skipping this image.")
                    except Exception as e:
                        print("Warning: An error occurred while processing an image: {}".format(str(e)))
            
            slides.append(slide_content)
        
        return slides